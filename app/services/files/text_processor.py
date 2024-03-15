import logging
import unicodedata
from io import BytesIO
import pandas as pd
from docx import Document
from pptx import Presentation  # Импортируем библиотеку python-pptx

import tiktoken
from unstructured.partition.auto import partition


class TextProcessor:
    CHARS_PER_TOKEN = 3
    PUNCTUATION = [".", "?", "!", "\n"]

    def __init__(self, encoding_name="cl100k_base"):
        self.encoding_name = encoding_name
        self.tokenizer = tiktoken.get_encoding(encoding_name)

    def chunk_texts(self, texts, min_tokens=50, max_tokens=500):
        if not isinstance(texts, str):
            texts = "\n".join(texts)

        text = self.normalize_unicode(texts)
        tokens = self.tokenizer.encode_ordinary(text)
        output = []
        i = 0
        while i < len(tokens):
            chunk_tokens = tokens[i: i + max_tokens]
            chunk = self.tokenizer.decode(chunk_tokens)
            last_punctuation = max([chunk.rfind(p) for p in self.PUNCTUATION], default=-1)
            if last_punctuation != -1 and last_punctuation > self.CHARS_PER_TOKEN * min_tokens:
                chunk = chunk[: last_punctuation + 1]

            i += len(self.tokenizer.encode_ordinary(chunk))
            output.append(chunk)
        return output

    @staticmethod
    def normalize_unicode(text):
        """Get rid of ligatures"""
        return unicodedata.normalize("NFKC", text)

    def extract_texts(self, data: bytes, file_type: str = "pdf") -> list[str]:
        file_like = BytesIO(data)
        try:
            if file_type == "pdf":
                elements = partition(file=file_like)
                texts = [element.text for element in elements]
            elif file_type in ["csv", "xlsx"]:
                df = pd.read_csv(file_like, engine='python') if file_type == "csv" else pd.read_excel(file_like)
                texts = df.applymap(str).apply(lambda x: ' '.join(x), axis=1).tolist()
            elif file_type == "txt":
                file_like.seek(0)
                texts = file_like.read().decode("utf-8").splitlines()
            elif file_type == "pptx":
                prs = Presentation(file_like)
                texts = [shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")]
            elif file_type == "docx":
                doc = Document(file_like)
                texts = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
            elif file_type == "doc":
                raise ValueError("Please convert your file to docx")
            elif file_type == "ppt":
                raise ValueError("Please convert your file to pptx")
            else:
                raise ValueError("Unsupported file type")
        except ValueError as ve:
            logging.error(f"Value Error: {str(ve)}")
            return []
        except Exception as e:
            logging.exception(f"An unexpected error occurred: {str(e)}")
            return []
        return texts
